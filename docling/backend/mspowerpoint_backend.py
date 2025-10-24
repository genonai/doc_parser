import logging
import re
import itertools
import os
import uuid
import shutil
import subprocess
from io import BytesIO
from pathlib import Path
from typing import Set, Union
from datetime import datetime, timedelta

from docling_core.types.doc.base import BoundingBox, CoordOrigin, Size
from docling_core.types.doc.labels import DocItemLabel, GroupLabel
from docling_core.types.doc.tokens import DocumentToken, TableToken, _LOC_PREFIX
from docling_core.types.doc.document import (
    ContentLayer,
    DoclingDocument,
    DocumentOrigin,
    DocTagsDocument,
    ImageRef,
    PictureClassificationClass,
    PictureClassificationData,
    PictureTabularChartData,
    ProvenanceItem,
    TableCell,
    TableData,
)
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.text import CT_TextLineBreak
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None  # type: ignore

from docling.backend.abstract_backend import (
    DeclarativeDocumentBackend,
    PaginatedDocumentBackend,
)
from docling.datamodel.base_models import InputFormat
from docling.datamodel.document import InputDocument

_log = logging.getLogger(__name__)


class MsPowerpointDocumentBackend(DeclarativeDocumentBackend, PaginatedDocumentBackend):
    def __init__(self, in_doc: "InputDocument", path_or_stream: Union[BytesIO, Path]):
        super().__init__(in_doc, path_or_stream)
        self.namespaces = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        # Powerpoint file:
        self.path_or_stream = path_or_stream

        self.pptx_obj = None
        self.valid = False
        self._exported_pdf_path: str | None = None
        # Chart handling mode: "image" (crop and save) or "serialize" (DocTags-based)
        try:
            mode_env = str(os.getenv("CHART_MODE", "image")).strip().lower()
        except Exception:
            mode_env = "image"
        self.chart_mode = mode_env if mode_env in {"image", "serialize"} else "image"
        try:
            if isinstance(self.path_or_stream, BytesIO):
                self.pptx_obj = Presentation(self.path_or_stream)
            elif isinstance(self.path_or_stream, Path):
                self.pptx_obj = Presentation(str(self.path_or_stream))

            self.valid = True
        except Exception as e:
            raise RuntimeError(
                f"MsPowerpointDocumentBackend could not load document with hash {self.document_hash}"
            ) from e

        return


    def _convert_current_pptx_to_pdf(self) -> str | None:
        """현재 PPTX를 LibreOffice로 PDF로 변환하고 경로를 반환. 실패 시 None."""
        try:
            if not isinstance(self.path_or_stream, Path):
                return None
            p: Path = self.path_or_stream
            if not p.exists():
                return None
            if self._exported_pdf_path and os.path.exists(self._exported_pdf_path):
                return self._exported_pdf_path
            soffice = shutil.which("soffice") or shutil.which("libreoffice")
            if not soffice:
                return None
            out_dir = os.getenv("CHART_IMG_DIR", "/workspaces/doc_parser/scratch/charts")
            try:
                os.makedirs(out_dir, exist_ok=True)
            except Exception:
                out_dir = "."
            cmd = [
                soffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", out_dir,
                str(p),
            ]
            try:
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            except Exception:
                return None
            pdf_candidate = Path(out_dir) / f"{p.stem}.pdf"
            if pdf_candidate.exists():
                self._exported_pdf_path = str(pdf_candidate)
                return self._exported_pdf_path
            return None
        except Exception:
            return None

    def _build_otsl_for_chart(self, categories, series_list) -> str:
        """Build OTSL string for a simple chart table (no spans).

        Layout:
        - Header row: <ched>Category + <ched>{series_name}...
        - Data rows: <rhed>{category} + value cells (<fcel>text or <ecel>)
        """
        parts: list[str] = []
        parts.append(f"<{DocumentToken.OTSL.value}>")

        # Header row
        parts.append(TableToken.OTSL_CHED.value)
        parts.append("Category")
        for s_name, _ in series_list:
            parts.append(TableToken.OTSL_CHED.value)
            parts.append("" if s_name is None else str(s_name))
        parts.append(TableToken.OTSL_NL.value)

        # If no series, add one generic Value column so we have at least 2 cols
        has_series = len(series_list) > 0
        if not has_series:
            series_list = [("Value", [])]

        # Data rows
        def _normalize_category_label(val):
            try:
                num = float(val)
                if 20000 <= num <= 60000:
                    base = datetime(1899, 12, 30)
                    return (base + timedelta(days=num)).strftime("%Y-%m-%d")
            except Exception:
                pass
            return "" if val is None else str(val)

        for i, cat in enumerate(categories):
            parts.append(TableToken.OTSL_RHED.value)
            parts.append(_normalize_category_label(cat))
            for _, values in series_list:
                val = values[i] if i < len(values) else None
                if val is None:
                    parts.append(TableToken.OTSL_ECEL.value)
                else:
                    parts.append(TableToken.OTSL_FCEL.value)
                    parts.append(str(val))
            parts.append(TableToken.OTSL_NL.value)

        parts.append(f"</{DocumentToken.OTSL.value}>")
        return "".join(parts)

    def _parse_table_content_from_otsl(self, otsl_content: str) -> TableData:
        """Parse OTSL content into TableData. Mirrors DocTags loader logic."""
        def otsl_extract_tokens_and_text(s: str):
            pattern = r"(<[^>]+>)"
            tokens = re.findall(pattern, s)
            tokens = [
                token
                for token in tokens
                if not (
                    token.startswith(f"<{_LOC_PREFIX}")
                    or token in [f"<{DocumentToken.OTSL.value}>", f"</{DocumentToken.OTSL.value}>"]
                )
            ]
            text_parts = re.split(pattern, s)
            text_parts = [
                token
                for token in text_parts
                if not (
                    token.startswith(f"<{_LOC_PREFIX}")
                    or token in [f"<{DocumentToken.OTSL.value}>", f"</{DocumentToken.OTSL.value}>"]
                )
            ]
            text_parts = [part for part in text_parts if str(part).strip()]
            return tokens, text_parts

        def otsl_parse_texts(texts, tokens):
            split_word = TableToken.OTSL_NL.value
            clean_tokens = []
            for t in tokens:
                if t in [
                    TableToken.OTSL_ECEL.value,
                    TableToken.OTSL_FCEL.value,
                    TableToken.OTSL_LCEL.value,
                    TableToken.OTSL_UCEL.value,
                    TableToken.OTSL_XCEL.value,
                    TableToken.OTSL_NL.value,
                    TableToken.OTSL_CHED.value,
                    TableToken.OTSL_RHED.value,
                    TableToken.OTSL_SROW.value,
                ]:
                    clean_tokens.append(t)
            tokens = clean_tokens
            split_row_tokens = [
                list(y)
                for x, y in itertools.groupby(tokens, lambda z: z == split_word)
                if not x
            ]

            table_cells = []
            r_idx = 0
            c_idx = 0

            def count_right(tokens, c_idx, r_idx, which_tokens):
                span = 0
                c_idx_iter = c_idx
                while tokens[r_idx][c_idx_iter] in which_tokens:
                    c_idx_iter += 1
                    span += 1
                    if c_idx_iter >= len(tokens[r_idx]):
                        return span
                return span

            def count_down(tokens, c_idx, r_idx, which_tokens):
                span = 0
                r_idx_iter = r_idx
                while tokens[r_idx_iter][c_idx] in which_tokens:
                    r_idx_iter += 1
                    span += 1
                    if r_idx_iter >= len(tokens):
                        return span
                return span

            for i, text in enumerate(texts):
                cell_text = ""
                if text in [
                    TableToken.OTSL_FCEL.value,
                    TableToken.OTSL_ECEL.value,
                    TableToken.OTSL_CHED.value,
                    TableToken.OTSL_RHED.value,
                    TableToken.OTSL_SROW.value,
                ]:
                    row_span = 1
                    col_span = 1
                    right_offset = 1
                    if text != TableToken.OTSL_ECEL.value:
                        cell_text = texts[i + 1]
                        right_offset = 2

                    next_right_cell = ""
                    if i + right_offset < len(texts):
                        next_right_cell = texts[i + right_offset]

                    next_bottom_cell = ""
                    if r_idx + 1 < len(split_row_tokens):
                        if c_idx < len(split_row_tokens[r_idx + 1]):
                            next_bottom_cell = split_row_tokens[r_idx + 1][c_idx]

                    if next_right_cell in [
                        TableToken.OTSL_LCEL.value,
                        TableToken.OTSL_XCEL.value,
                    ]:
                        col_span += count_right(
                            split_row_tokens,
                            c_idx + 1,
                            r_idx,
                            [TableToken.OTSL_LCEL.value, TableToken.OTSL_XCEL.value],
                        )
                    if next_bottom_cell in [
                        TableToken.OTSL_UCEL.value,
                        TableToken.OTSL_XCEL.value,
                    ]:
                        row_span += count_down(
                            split_row_tokens,
                            c_idx,
                            r_idx + 1,
                            [TableToken.OTSL_UCEL.value, TableToken.OTSL_XCEL.value],
                        )

                    table_cells.append(
                        TableCell(
                            text=str(cell_text).strip(),
                            row_span=row_span,
                            col_span=col_span,
                            start_row_offset_idx=r_idx,
                            end_row_offset_idx=r_idx + row_span,
                            start_col_offset_idx=c_idx,
                            end_col_offset_idx=c_idx + col_span,
                        )
                    )
                if text in [
                    TableToken.OTSL_FCEL.value,
                    TableToken.OTSL_ECEL.value,
                    TableToken.OTSL_CHED.value,
                    TableToken.OTSL_RHED.value,
                    TableToken.OTSL_SROW.value,
                    TableToken.OTSL_LCEL.value,
                    TableToken.OTSL_UCEL.value,
                    TableToken.OTSL_XCEL.value,
                ]:
                    c_idx += 1
                if text == TableToken.OTSL_NL.value:
                    r_idx += 1
                    c_idx = 0
            return table_cells, split_row_tokens

        tokens, mixed_texts = otsl_extract_tokens_and_text(otsl_content)
        table_cells, split_row_tokens = otsl_parse_texts(mixed_texts, tokens)
        return TableData(
            num_rows=len(split_row_tokens),
            num_cols=(max(len(row) for row in split_row_tokens) if split_row_tokens else 0),
            table_cells=table_cells,
        )

    def _maybe_format_categories_with_xml(self, chart, categories):
        """카테고리 값이 Excel 날짜 직렬값일 때, 차트 XML의 formatCode를 이용해 사람이 읽을 수 있는 날짜 문자열로 변환."""
        if not categories:
            return categories
        try:
            elem = getattr(chart, "_element", None)
            if elem is None:
                return categories
            # 1) 시리즈 카테고리 캐시의 포맷 코드
            fmt_nodes = elem.xpath(
                ".//c:cat//c:numCache/c:formatCode",
                namespaces=self.namespaces,
            )
            fmt = fmt_nodes[0].text if fmt_nodes else None
            # 2) 축 포맷 코드 보조 경로
            if not fmt:
                fmt_attr = elem.xpath(
                    ".//c:catAx/c:numFmt/@formatCode",
                    namespaces=self.namespaces,
                )
                fmt = fmt_attr[0] if fmt_attr else None

            def excel_serial_to_date_str(val, format_code: str | None) -> str:
                try:
                    num = float(val)
                except Exception:
                    return str(val)
                # 휴리스틱: Excel date 범위
                if not (20000 <= num <= 60000):
                    return str(val)
                base = datetime(1899, 12, 30)
                dt = base + timedelta(days=num)
                pfmt = "%Y-%m-%d"
                if format_code:
                    fl = format_code.lower()
                    # 간단 매핑 (필요 시 확장)
                    if fl in ("m/d/yyyy", "mm/dd/yyyy"):
                        pfmt = "%m/%d/%Y"
                    elif fl in ("yyyy-mm-dd",):
                        pfmt = "%Y-%m-%d"
                    elif ("y" in fl and "m" in fl and "d" in fl):
                        pfmt = "%Y-%m-%d"
                return dt.strftime(pfmt)

            if fmt:
                return [excel_serial_to_date_str(c, fmt) for c in categories]
            else:
                # 포맷 코드가 없을 때도 직렬값이면 ISO로 변환
                return [excel_serial_to_date_str(c, None) for c in categories]
        except Exception:
            return categories

    def _extract_categories_from_xml(self, chart):
        """차트 XML에서 c:cat의 캐시(strCache/numCache)를 직접 읽어 카테고리를 추출.
        - 날짜인 경우 formatCode에 맞춰 문자열로 변환
        - 문자열 캐시가 있으면 그대로 사용
        """
        try:
            elem = getattr(chart, "_element", None)
            if elem is None:
                return []
            # 우선 strCache
            str_nodes = elem.xpath(
                ".//c:cat//c:strCache/c:pt/c:v",
                namespaces=self.namespaces,
            )
            if str_nodes:
                return [str(n.text) if n is not None else "" for n in str_nodes]

            # 숫자 캐시 + 포맷 코드
            fmt_nodes = elem.xpath(
                ".//c:cat//c:numCache/c:formatCode",
                namespaces=self.namespaces,
            )
            fmt = fmt_nodes[0].text if fmt_nodes else None
            num_nodes = elem.xpath(
                ".//c:cat//c:numCache/c:pt/c:v",
                namespaces=self.namespaces,
            )
            if not num_nodes:
                return []

            def excel_serial_to_date_str(val, format_code: str | None) -> str:
                try:
                    num = float(val)
                except Exception:
                    return str(val)
                if not (20000 <= num <= 60000):
                    return str(val)
                base = datetime(1899, 12, 30)
                dt = base + timedelta(days=num)
                pfmt = "%Y-%m-%d"
                if format_code:
                    fl = format_code.lower()
                    if fl in ("m/d/yyyy", "mm/dd/yyyy"):
                        pfmt = "%m/%d/%Y"
                    elif fl in ("yyyy-mm-dd",):
                        pfmt = "%Y-%m-%d"
                    elif ("y" in fl and "m" in fl and "d" in fl):
                        pfmt = "%Y-%m-%d"
                return dt.strftime(pfmt)

            values = [n.text for n in num_nodes]
            return [excel_serial_to_date_str(v, fmt) for v in values]
        except Exception:
            return []

    def page_count(self) -> int:
        if self.is_valid():
            assert self.pptx_obj is not None
            return len(self.pptx_obj.slides)
        else:
            return 0

    def is_valid(self) -> bool:
        return self.valid

    @classmethod
    def supports_pagination(cls) -> bool:
        return True  # True? if so, how to handle pages...

    def unload(self):
        if isinstance(self.path_or_stream, BytesIO):
            self.path_or_stream.close()

        self.path_or_stream = None

    @classmethod
    def supported_formats(cls) -> Set[InputFormat]:
        return {InputFormat.PPTX}

    def convert(self) -> DoclingDocument:
        # Parses the PPTX into a structured document model.
        # origin = DocumentOrigin(filename=self.path_or_stream.name, mimetype=next(iter(FormatToMimeType.get(InputFormat.PPTX))), binary_hash=self.document_hash)

        # # binary_hash는 Uint64 범위여야 하므로 64비트 마스킹 적용
        # if isinstance(self.document_hash, str):
        #     try:
        #         bh = int(self.document_hash, 16) & 0xFFFFFFFFFFFFFFFF
        #     except Exception:
        #         bh = 0
        # elif isinstance(self.document_hash, int):
        #     bh = self.document_hash & 0xFFFFFFFFFFFFFFFF
        # else:
        #     bh = 0

        origin = DocumentOrigin(
            filename=self.file.name or "file",
            mimetype="application/vnd.ms-powerpoint",
            binary_hash=0 if isinstance(self.document_hash, str) else (self.document_hash & 0xFFFFFFFFFFFFFFFF),
            # binary_hash=bh,
        )

        doc = DoclingDocument(
            name=self.file.stem or "file", origin=origin
        )  # must add origin information
        doc = self.walk_linear(self.pptx_obj, doc)

        return doc

    def generate_prov(
        self, shape, slide_ind, text="", slide_size=Size(width=1, height=1)
    ):
        if shape.left:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
        else:
            left = 0
            top = 0
            width = slide_size.width
            height = slide_size.height
        shape_bbox = (left, top, left + width, top + height)
        shape_bbox = BoundingBox.from_tuple(shape_bbox, origin=CoordOrigin.BOTTOMLEFT)
        prov = ProvenanceItem(
            page_no=slide_ind + 1, charspan=(0, len(text)), bbox=shape_bbox
        )

        return prov

    def handle_text_elements(
        self, shape, parent_slide, slide_ind, doc: DoclingDocument, slide_size
    ):
        is_list_group_created = False
        enum_list_item_value = 0
        new_list = None
        doc_label = DocItemLabel.LIST_ITEM
        prov = self.generate_prov(shape, slide_ind, shape.text.strip(), slide_size)

        def is_list_item(paragraph):
            """Check if the paragraph is a list item."""
            p = paragraph._element
            if (
                p.find(".//a:buChar", namespaces={"a": self.namespaces["a"]})
                is not None
            ):
                return (True, "Bullet")
            elif (
                p.find(".//a:buAutoNum", namespaces={"a": self.namespaces["a"]})
                is not None
            ):
                return (True, "Numbered")
            elif paragraph.level > 0:
                # Most likely a sub-list
                return (True, "None")
            else:
                return (False, "None")

        # Iterate through paragraphs to build up text
        for paragraph in shape.text_frame.paragraphs:
            is_a_list, bullet_type = is_list_item(paragraph)
            p = paragraph._element

            # Convert line breaks to spaces and accumulate text
            p_text = ""
            for e in p.content_children:
                if isinstance(e, CT_TextLineBreak):
                    p_text += " "
                else:
                    p_text += e.text

            if is_a_list:
                enum_marker = ""
                enumerated = bullet_type == "Numbered"

                if not is_list_group_created:
                    new_list = doc.add_list_group(
                        name="list",
                        parent=parent_slide,
                    )
                    is_list_group_created = True
                    enum_list_item_value = 0

                if enumerated:
                    enum_list_item_value += 1
                    enum_marker = str(enum_list_item_value) + "."

                doc.add_list_item(
                    marker=enum_marker,
                    enumerated=enumerated,
                    parent=new_list,
                    text=p_text,
                    prov=prov,
                )
            else:  # is paragraph not a list item
                # Assign proper label to the text, depending if it's a Title or Section Header
                # For other types of text, assign - PARAGRAPH
                # doc_label = DocItemLabel.PARAGRAPH
                doc_label = DocItemLabel.TEXT # text 로 변경(toc enrichment 적용 위함)
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type in [
                        PP_PLACEHOLDER.CENTER_TITLE,
                        PP_PLACEHOLDER.TITLE,
                    ]:
                        # It's a title
                        doc_label = DocItemLabel.TITLE
                    elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                        DocItemLabel.SECTION_HEADER

                # Caption detection: starts with "자료:" or "참고:"
                stripped = p_text.lstrip()
                if stripped.startswith("자료:") or stripped.startswith("참고:") or stripped.startswith("출처:"):
                    doc_label = DocItemLabel.CAPTION
                # if stripped.startswith("*"):
                #     doc_label = DocItemLabel.FOOTNOTE

                # output accumulated inline text:
                doc.add_text(
                    label=doc_label,
                    parent=parent_slide,
                    text=p_text,
                    prov=prov,
                )
        return

    def handle_title(self, shape, parent_slide, slide_ind, doc):
        placeholder_type = shape.placeholder_format.type
        txt = shape.text.strip()
        prov = self.generate_prov(shape, slide_ind, txt)

        if len(txt.strip()) > 0:
            # title = slide.shapes.title.text if slide.shapes.title else "No title"
            if placeholder_type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE]:
                _log.info(f"Title found: {shape.text}")
                doc.add_text(
                    label=DocItemLabel.TITLE, parent=parent_slide, text=txt, prov=prov
                )
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                _log.info(f"Subtitle found: {shape.text}")
                # Using DocItemLabel.FOOTNOTE, while SUBTITLE label is not avail.
                doc.add_text(
                    label=DocItemLabel.SECTION_HEADER,
                    parent=parent_slide,
                    text=txt,
                    prov=prov,
                )
        return

    def handle_pictures(self, shape, parent_slide, slide_ind, doc, slide_size):
        # Open it with PIL
        if not hasattr(shape, "image") or shape.image is None:
            return
        try:
            # Get the image bytes
            image = shape.image
            image_bytes = image.blob
            im_dpi, _ = image.dpi
            pil_image = Image.open(BytesIO(image_bytes))

            # shape has picture
            prov = self.generate_prov(shape, slide_ind, "", slide_size)
            doc.add_picture(
                parent=parent_slide,
                image=ImageRef.from_pil(image=pil_image, dpi=im_dpi),
                caption=None,
                prov=prov,
            )
        except (UnidentifiedImageError, OSError) as e:
            _log.warning(f"Warning: image cannot be loaded by Pillow: {e}")
        return
    
    def handle_charts(self, shape, parent_slide, slide_ind, doc, slide_size):
        """차트를 직렬화(OTSL/DocTags) 기반으로 처리하여 annotation을 포함해 Picture로 추가."""
        if not hasattr(shape, "has_chart") or not shape.has_chart:
            return

        prov = self.generate_prov(shape, slide_ind, "", slide_size)

        # 시리즈/카테고리 수집
        series_list = []  # [(name, [values...])]
        max_points = 0
        categories = []
        try:
            chart = shape.chart
        except Exception:
            chart = None

        if chart is not None:
            for s in getattr(chart, "series", []) or []:
                try:
                    s_name = s.name if hasattr(s, "name") else "Series"
                except Exception:
                    s_name = "Series"
                values = []
                try:
                    pts = getattr(s, "points", [])
                    if pts:
                        for pt in pts:
                            try:
                                values.append(pt.value)
                            except Exception:
                                values.append(None)
                    else:
                        for v in getattr(s, "values", []) or []:
                            values.append(v)
                except Exception:
                    pass
                max_points = max(max_points, len(values))
                series_list.append((s_name, values))

            # 카테고리는 OOXML에서 직접 우선 추출 (formatCode 반영)
            categories = self._extract_categories_from_xml(chart)
            if not categories:
                try:
                    plots = getattr(chart, "plots", [])
                    if plots and hasattr(plots[0], "categories") and plots[0].categories:
                        for c in plots[0].categories:
                            try:
                                categories.append(str(c))
                            except Exception:
                                categories.append("")
                except Exception:
                    categories = []

        if not categories and max_points > 0:
            categories = [f"idx_{i+1}" for i in range(max_points)]

        # 카테고리 날짜 포맷 보정 (보조 휴리스틱)
        categories = self._maybe_format_categories_with_xml(chart, categories)

        # OTSL 구성 및 파싱
        otsl_text = self._build_otsl_for_chart(categories=categories, series_list=series_list)

        # 차트 타입 태그 결정
        chart_type_name = "tabular_chart"
        try:
            from pptx.enum.chart import XL_CHART_TYPE
            ct = getattr(chart, "chart_type", None)
            if ct is not None:
                if ct in {
                    getattr(XL_CHART_TYPE, "LINE", None),
                    getattr(XL_CHART_TYPE, "LINE_MARKERS", None),
                    getattr(XL_CHART_TYPE, "LINE_MARKERS_STACKED", None),
                    getattr(XL_CHART_TYPE, "LINE_STACKED", None),
                }:
                    chart_type_name = "line_chart"
                elif ct in {
                    getattr(XL_CHART_TYPE, "BAR_CLUSTERED", None),
                    getattr(XL_CHART_TYPE, "BAR_STACKED", None),
                }:
                    chart_type_name = "bar_chart"
                elif ct in {
                    getattr(XL_CHART_TYPE, "COLUMN_CLUSTERED", None),
                    getattr(XL_CHART_TYPE, "COLUMN_STACKED", None),
                    getattr(XL_CHART_TYPE, "COLUMN_STACKED_100", None),
                }:
                    chart_type_name = "bar_chart"
                elif ct in {
                    getattr(XL_CHART_TYPE, "BAR_STACKED_100", None),
                }:
                    chart_type_name = "stacked_bar_chart"
                elif ct in {
                    getattr(XL_CHART_TYPE, "PIE", None),
                    getattr(XL_CHART_TYPE, "PIE_EXPLODED", None),
                }:
                    chart_type_name = "pie_chart"
                elif ct in {
                    getattr(XL_CHART_TYPE, "XY_SCATTER", None),
                    getattr(XL_CHART_TYPE, "XY_SCATTER_LINES", None),
                    getattr(XL_CHART_TYPE, "XY_SCATTER_LINES_NO_MARKERS", None),
                }:
                    chart_type_name = "scatter_chart"
        except Exception:
            pass

        # 위치 토큰 생성
        try:
            page_w = float(slide_size.width) if slide_size and slide_size.width else 1.0
            page_h = float(slide_size.height) if slide_size and slide_size.height else 1.0
            left = float(getattr(shape, "left", 0) or 0)
            top = float(getattr(shape, "top", 0) or 0)
            width = float(getattr(shape, "width", 1) or 1)
            height = float(getattr(shape, "height", 1) or 1)
            l_norm = max(0.0, min(1.0, left / page_w))
            t_norm = max(0.0, min(1.0, top / page_h))
            r_norm = max(0.0, min(1.0, (left + width) / page_w))
            b_norm = max(0.0, min(1.0, (top + height) / page_h))
            loc_tokens = (
                f"<{_LOC_PREFIX}{int(l_norm * 500)}>"
                f"<{_LOC_PREFIX}{int(t_norm * 500)}>"
                f"<{_LOC_PREFIX}{int(r_norm * 500)}>"
                f"<{_LOC_PREFIX}{int(b_norm * 500)}>"
            )
        except Exception:
            loc_tokens = ""

        # DocTags 페이지 구성 및 로드
        if chart_type_name != "tabular_chart":
            chart_type_tag_open = f"<{chart_type_name}>"
            chart_type_tag_close = f"</{chart_type_name}>"
        else:
            chart_type_tag_open = ""
            chart_type_tag_close = ""

        doctags_page = (
            f"<{DocumentToken.CHART.value}>"
            f"{loc_tokens}"
            f"{chart_type_tag_open}"
            f"{otsl_text}"
            f"{chart_type_tag_close}"
            f"</{DocumentToken.CHART.value}>"
        )

        annotations = []
        try:
            dt_doc = DocTagsDocument.from_doctags_and_image_pairs([doctags_page], images=None)
            tmp_doc = DoclingDocument.load_from_doctags(dt_doc, document_name="tmp")
            tmp_pic = next(iter(tmp_doc.pictures), None)
            if tmp_pic is not None:
                for ann in tmp_pic.annotations:
                    annotations.append(ann)
        except Exception:
            chart_table = self._parse_table_content_from_otsl(otsl_text)
            annotations.append(
                PictureClassificationData(
                    provenance="load_from_doctags",
                    predicted_classes=[
                        PictureClassificationClass(
                            class_name=chart_type_name, confidence=1.0
                        )
                    ],
                )
            )
            annotations.append(
                PictureTabularChartData(
                    chart_data=chart_table,
                    title=getattr(shape, "name", "chart") or "chart",
                )
            )

        # 투명 placeholder 이미지 부착 (직렬화 전용, 실이미지 렌더링 없음)
        try:
            placeholder_img = Image.new("RGBA", (2, 2), (0, 0, 0, 0))
            placeholder_ref = ImageRef.from_pil(image=placeholder_img, dpi=72)
        except Exception:
            placeholder_ref = None

        pic = doc.add_picture(
            parent=parent_slide,
            prov=prov,
            image=placeholder_ref,
            annotations=annotations,
        )
        try:
            pic.label = DocItemLabel.CHART
        except Exception:
            pass
        return

    def handle_charts_save_image(self, shape, parent_slide, slide_ind, doc, slide_size):
        """차트를 PDF 렌더링 후 해당 영역만 크롭하여 이미지로 저장하고 Picture로 추가."""
        if not hasattr(shape, "has_chart") or not shape.has_chart:
            return

        prov = self.generate_prov(shape, slide_ind, "", slide_size)

        # 차트를 PDF로 변환 후 해당 영역 크롭하여 이미지 저장
        image_ref = None
        try:
            pdf_path = self._convert_current_pptx_to_pdf()
            if fitz is None or not pdf_path or not os.path.exists(pdf_path):
                raise RuntimeError("PDF export unavailable")

            dpi = int(os.getenv("CHART_PDF_DPI", "144"))
            zoom = dpi / 72.0
            doc_pdf = fitz.open(pdf_path)
            page_index = max(0, min(slide_ind, len(doc_pdf) - 1))
            page = doc_pdf[page_index]
            mat = fitz.Matrix(zoom, zoom)
            get_pix = getattr(page, "get_pixmap", None)
            if get_pix is None:
                get_pix = getattr(page, "getPixmap", None)
            if get_pix is None:
                raise RuntimeError("PyMuPDF get_pixmap not available")
            pix = get_pix(matrix=mat, alpha=False)

            # 슬라이드 좌표(EMU) -> 슬라이드 크기 기준 정규화 -> 픽셀 좌표 변환
            left = float(getattr(shape, "left", 0) or 0)
            top = float(getattr(shape, "top", 0) or 0)
            width = float(getattr(shape, "width", 1) or 1)
            height = float(getattr(shape, "height", 1) or 1)
            slide_w = float(slide_size.width) if slide_size and slide_size.width else 1.0
            slide_h = float(slide_size.height) if slide_size and slide_size.height else 1.0

            l_norm = max(0.0, min(1.0, left / slide_w))
            t_norm = max(0.0, min(1.0, top / slide_h))
            r_norm = max(0.0, min(1.0, (left + width) / slide_w))
            b_norm = max(0.0, min(1.0, (top + height) / slide_h))

            # 픽셀 좌표 (이미지 좌표계는 상단 원점)
            l_px = int(round(l_norm * pix.width))
            r_px = int(round(r_norm * pix.width))
            t_px = int(round(t_norm * pix.height))
            b_px = int(round(b_norm * pix.height))

            # 좌표 정렬 및 최소 크기 검증
            x0, x1 = sorted((max(0, min(l_px, pix.width)), max(0, min(r_px, pix.width))))
            y0, y1 = sorted((max(0, min(t_px, pix.height)), max(0, min(b_px, pix.height))))
            if (x1 - x0) < 1 or (y1 - y0) < 1:
                raise RuntimeError("invalid crop box size")

            from PIL import Image as PILImage
            img_full = PILImage.frombytes("RGB", (pix.width, pix.height), pix.samples)
            crop_box = (x0, y0, x1, y1)
            cropped = img_full.crop(crop_box)

            # 크롭한 이미지를 ImageRef로 변환
            image_ref = ImageRef.from_pil(image=cropped, dpi=72)
            doc_pdf.close()
        except Exception as e:
            _log.debug(f"Chart image render failed: {e}")
            try:
                # 투명 플레이스홀더(차트 표식 유지)
                from PIL import Image
                ph = Image.new("RGBA", (2, 2), (0, 0, 0, 0))
                image_ref = ImageRef.from_pil(image=ph, dpi=72)
            except Exception:
                pass

        if image_ref is not None:
            # add_picture으로 전달
            doc.add_picture(
                parent=parent_slide,
                image=image_ref,  # 크롭한 이미지의 ImageRef 사용
                caption=None,
                prov=prov,
            )

        return
    
    def handle_tables(self, shape, parent_slide, slide_ind, doc, slide_size):
        # Handling tables, images, charts
        if shape.has_table:
            table = shape.table
            table_xml = shape._element

            prov = self.generate_prov(shape, slide_ind, "", slide_size)

            num_cols = 0
            num_rows = len(table.rows)
            tcells = []
            # Access the XML element for the shape that contains the table
            table_xml = shape._element

            for row_idx, row in enumerate(table.rows):
                if len(row.cells) > num_cols:
                    num_cols = len(row.cells)
                for col_idx, cell in enumerate(row.cells):
                    # Access the XML of the cell (this is the 'tc' element in table XML)
                    cell_xml = table_xml.xpath(
                        f".//a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]"
                    )

                    if not cell_xml:
                        continue  # If no cell XML is found, skip

                    cell_xml = cell_xml[0]  # Get the first matching XML node
                    row_span = cell_xml.get("rowSpan")  # Vertical span
                    col_span = cell_xml.get("gridSpan")  # Horizontal span

                    if row_span is None:
                        row_span = 1
                    else:
                        row_span = int(row_span)

                    if col_span is None:
                        col_span = 1
                    else:
                        col_span = int(col_span)

                    icell = TableCell(
                        text=cell.text.strip(),
                        row_span=row_span,
                        col_span=col_span,
                        start_row_offset_idx=row_idx,
                        end_row_offset_idx=row_idx + row_span,
                        start_col_offset_idx=col_idx,
                        end_col_offset_idx=col_idx + col_span,
                        column_header=row_idx == 0,
                        row_header=False,
                    )
                    if len(cell.text.strip()) > 0:
                        tcells.append(icell)
            # Initialize Docling TableData
            data = TableData(num_rows=num_rows, num_cols=num_cols, table_cells=[])
            # Populate
            for tcell in tcells:
                data.table_cells.append(tcell)
            if len(tcells) > 0:
                # If table is not fully empty...
                # Create Docling table
                doc.add_table(parent=parent_slide, data=data, prov=prov)
        return

    def walk_linear(self, pptx_obj, doc) -> DoclingDocument:
        # Units of size in PPTX by default are EMU units (English Metric Units)
        slide_width = pptx_obj.slide_width
        slide_height = pptx_obj.slide_height

        max_levels = 10
        parents = {}  # type: ignore
        for i in range(max_levels):
            parents[i] = None

        # Loop through each slide
        for slide_num, slide in enumerate(pptx_obj.slides):
            slide_ind = pptx_obj.slides.index(slide)
            parent_slide = doc.add_group(
                name=f"slide-{slide_ind}", label=GroupLabel.CHAPTER, parent=parents[0]
            )

            slide_size = Size(width=slide_width, height=slide_height)
            doc.add_page(page_no=slide_ind + 1, size=slide_size)

            def handle_shapes(shape, parent_slide, slide_ind, doc, slide_size):
                handle_groups(shape, parent_slide, slide_ind, doc, slide_size)
                if shape.has_table:
                    # Handle Tables
                    self.handle_tables(shape, parent_slide, slide_ind, doc, slide_size)
                if hasattr(shape, "has_chart") and shape.has_chart:
                    if self.chart_mode == "serialize":
                        self.handle_charts(shape, parent_slide, slide_ind, doc, slide_size)
                    else:
                        self.handle_charts_save_image(shape, parent_slide, slide_ind, doc, slide_size)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Handle Pictures
                    if hasattr(shape, "image") and shape.image is not None:
                        self.handle_pictures(
                            shape, parent_slide, slide_ind, doc, slide_size
                        )
                # If shape doesn't have any text, move on to the next shape
                if not hasattr(shape, "text"):
                    return
                if shape.text is None:
                    return
                if len(shape.text.strip()) == 0:
                    return
                if not shape.has_text_frame:
                    _log.warning("Warning: shape has text but not text_frame")
                    return
                # Handle other text elements, including lists (bullet lists, numbered lists)
                self.handle_text_elements(
                    shape, parent_slide, slide_ind, doc, slide_size
                )
                return

            def handle_groups(shape, parent_slide, slide_ind, doc, slide_size):
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # 그룹 노드 생성하여 문맥 보존
                    group_parent = doc.add_group(
                        label=GroupLabel.INLINE,
                        name=getattr(shape, "name", None) or f"group-{slide_ind}-{getattr(shape, 'shape_id', 'id')}",
                        parent=parent_slide,
                    )

                    # 그룹의 children 좌표계 정보 (chOff/chExt) 추출
                    # NOTE: 이전 버전에서는 chOff/chExt를 사용해 자식 좌표를 절대좌표로 변환했습니다.
                    # 여기서는 변환을 적용하지 않고 로컬(top,left) 기준 정렬만 수행합니다.

                    # 그룹의 배치/크기 (부모 좌표계) — 사용하지 않음
                    g_left = getattr(shape, "left", 0) or 0
                    g_top = getattr(shape, "top", 0) or 0
                    g_width = getattr(shape, "width", 1) or 1
                    g_height = getattr(shape, "height", 1) or 1

                    # 변환 파라미터를 사용하지 않음

                    # 자식들을 시각적 순서(위→아래, 좌→우) 단순 정렬: 로컬 좌표 기준
                    def _abs_lt(s):
                        s_left = getattr(s, "left", 0) or 0
                        s_top = getattr(s, "top", 0) or 0
                        return (s_top, s_left)

                    sorted_children = sorted(list(shape.shapes), key=_abs_lt)
                    for groupedshape in sorted_children:
                        handle_shapes(
                            groupedshape, group_parent, slide_ind, doc, slide_size
                        )

            # Loop through each shape in the slide
            for shape in slide.shapes:
                handle_shapes(shape, parent_slide, slide_ind, doc, slide_size)

            # Handle notes slide
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame is not None:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        bbox = BoundingBox(l=0, t=0, r=0, b=0)
                        prov = ProvenanceItem(
                            page_no=slide_ind + 1,
                            charspan=(0, len(notes_text)),
                            bbox=bbox,
                        )
                        doc.add_text(
                            label=DocItemLabel.TEXT,
                            parent=parent_slide,
                            text=notes_text,
                            prov=prov,
                            content_layer=ContentLayer.FURNITURE,
                        )
                        
        return doc
